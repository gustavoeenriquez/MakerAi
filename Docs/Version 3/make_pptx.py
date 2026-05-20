"""
Script to generate MakerAI Capability System PPTX presentation.
Run: python make_pptx.py
Output: uMakerAi-CapabilitySystem.EN.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # slide background (deep navy)
C_ACCENT  = RGBColor(0x16, 0x21, 0x3E)   # card / box background
C_BLUE    = RGBColor(0x0F, 0x3D, 0x91)   # primary blue
C_TEAL    = RGBColor(0x00, 0xB4, 0xD8)   # teal highlight
C_ORANGE  = RGBColor(0xFF, 0x8C, 0x00)   # orange accent
C_GREEN   = RGBColor(0x2E, 0xCC, 0x71)   # green (gap / bridge)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY   = RGBColor(0xC8, 0xD6, 0xE5)   # light gray text
C_YELLOW  = RGBColor(0xFF, 0xD7, 0x00)   # yellow highlight

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bullet_box(slide, title, bullets,
                   l, t, w, h,
                   box_color=C_ACCENT, title_color=C_TEAL,
                   bullet_color=C_WHITE, title_size=16, bullet_size=13,
                   bullet_char="▸ "):
    """Draws a rounded-corner-ish box with title + bullet list."""
    add_rect(slide, l, t, w, h, fill=box_color,
             line=C_TEAL, line_w=Pt(1))
    # title
    add_text(slide, title, l + Inches(0.15), t + Inches(0.1),
             w - Inches(0.3), Inches(0.4),
             size=title_size, bold=True, color=title_color)
    # bullets
    txb = slide.shapes.add_textbox(
        l + Inches(0.15), t + Inches(0.52),
        w - Inches(0.3), h - Inches(0.65))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = bullet_char + b
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = bullet_color


def add_code_box(slide, code_lines, l, t, w, h, size=11):
    add_rect(slide, l, t, w, h, fill=RGBColor(0x0D, 0x0D, 0x1A),
             line=RGBColor(0x33, 0x33, 0x66), line_w=Pt(1))
    txb = slide.shapes.add_textbox(
        l + Inches(0.12), t + Inches(0.1),
        w - Inches(0.24), h - Inches(0.2))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.name = "Consolas"
        run.font.color.rgb = RGBColor(0xA8, 0xFF, 0x78)


def slide_title_bar(slide, title, subtitle=None):
    """Top bar with gradient-like colored stripe."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), fill=C_BLUE)
    add_text(slide, title,
             Inches(0.35), Inches(0.08), Inches(12), Inches(0.6),
             size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.65), Inches(12), Inches(0.38),
                 size=14, bold=False, color=C_TEAL)


def slide_section_divider(slide, number, title, subtitle=None):
    """Full-screen section title slide."""
    set_bg(slide, C_DARK)
    # accent bar left
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill=C_TEAL)
    # number
    add_text(slide, number,
             Inches(0.5), Inches(2.2), Inches(3), Inches(1.2),
             size=80, bold=True, color=RGBColor(0x1E, 0x3A, 0x6E))
    # title
    add_text(slide, title,
             Inches(0.5), Inches(3.2), Inches(12.3), Inches(1.3),
             size=40, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.7),
                 size=20, color=C_TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.18), fill=C_TEAL)
add_rect(sl, 0, SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18), fill=C_TEAL)
add_rect(sl, 0, Inches(0.18), Inches(0.18), SLIDE_H - Inches(0.36), fill=C_BLUE)

add_text(sl, "MakerAI v3.3",
         Inches(0.5), Inches(1.0), Inches(12), Inches(0.7),
         size=22, bold=False, color=C_TEAL)

add_text(sl, "Capability System",
         Inches(0.5), Inches(1.65), Inches(12), Inches(1.3),
         size=52, bold=True, color=C_WHITE)

add_text(sl, "ModelCaps & SessionCaps",
         Inches(0.5), Inches(2.85), Inches(12), Inches(0.8),
         size=34, bold=False, color=C_ORANGE)

add_text(sl,
         "How MakerAI automatically orchestrates LLM providers,\n"
         "bridges, and specialized endpoints through a unified capability model.",
         Inches(0.5), Inches(3.65), Inches(10), Inches(1.1),
         size=16, color=C_LGRAY)

add_rect(sl, Inches(0.5), Inches(5.0), Inches(4), Inches(0.06),
         fill=C_TEAL)

add_text(sl, "makerai.cimamaker.com  ·  March 2026",
         Inches(0.5), Inches(5.15), Inches(10), Inches(0.4),
         size=13, color=C_LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Agenda", "What we'll cover in this session")

items = [
    ("1", "Introduction — why a unified capability model?"),
    ("2", "TAiCapability enum — the vocabulary"),
    ("3", "ModelCaps vs SessionCaps — the two key properties"),
    ("4", "Gap Analysis — the orchestration engine"),
    ("5", "Configuration via TAiChatFactory"),
    ("6", "Runtime configuration in code"),
    ("7", "Backward compatibility with legacy params"),
    ("8", "10 common configuration patterns"),
    ("9", "Provider capability reference"),
    ("10", "Adding a new provider"),
]

col_w = Inches(5.9)
for i, (num, text) in enumerate(items):
    col = i % 2
    row = i // 2
    lx = Inches(0.4) + col * (col_w + Inches(0.5))
    ty = Inches(1.25) + row * Inches(1.05)
    add_rect(sl, lx, ty, col_w, Inches(0.88),
             fill=C_ACCENT, line=C_BLUE, line_w=Pt(1))
    add_text(sl, num, lx + Inches(0.12), ty + Inches(0.14),
             Inches(0.5), Inches(0.55),
             size=22, bold=True, color=C_TEAL)
    add_text(sl, text,
             lx + Inches(0.55), ty + Inches(0.18),
             col_w - Inches(0.7), Inches(0.55),
             size=14, color=C_WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Introduction: the problem
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Introduction", "The problem with the legacy system")

add_bullet_box(sl,
    "Legacy: 4 scattered parameters",
    [
        "NativeInputFiles  — physical file types the model accepts",
        "NativeOutputFiles — file types the model generates",
        "ChatMediaSupports — logical native capabilities",
        "EnabledFeatures   — desired capabilities for the session",
    ],
    Inches(0.35), Inches(1.25), Inches(5.9), Inches(2.5),
    box_color=RGBColor(0x2D, 0x10, 0x10),
    title_color=C_ORANGE,
    bullet_color=C_LGRAY,
    bullet_size=13)

add_bullet_box(sl,
    "New in v3.3: 2 unified properties",
    [
        "ModelCaps   — what the model can do natively via completions",
        "SessionCaps — what this session needs (may exceed native caps)",
        "Gap = SessionCaps − ModelCaps  →  auto-selects the right bridge",
        "Backward-compatible: legacy models keep working unchanged",
    ],
    Inches(6.7), Inches(1.25), Inches(6.25), Inches(2.5),
    box_color=RGBColor(0x0A, 0x2A, 0x1A),
    title_color=C_GREEN,
    bullet_size=13)

# arrow between boxes
add_rect(sl, Inches(6.25), Inches(2.3), Inches(0.45), Inches(0.06),
         fill=C_TEAL)
add_text(sl, "→", Inches(6.28), Inches(2.1), Inches(0.5), Inches(0.4),
         size=22, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

add_text(sl, "One pair of lines now fully describes any model:",
         Inches(0.35), Inches(3.9), Inches(12.5), Inches(0.4),
         size=14, bold=True, color=C_LGRAY)

add_code_box(sl, [
    "TAiChatFactory.Instance.RegisterUserParam('OpenAi', 'gpt-4.1', 'ModelCaps',   '[cap_Image]');",
    "TAiChatFactory.Instance.RegisterUserParam('OpenAi', 'gpt-4.1', 'SessionCaps', '[cap_Image]');",
], Inches(0.35), Inches(4.35), Inches(12.6), Inches(0.95))

add_text(sl,
    "The orchestrator reads the Gap and automatically activates transcription bridges, "
    "image-description bridges, TTS endpoints, image-generation endpoints, and more — "
    "with zero extra code in your application.",
    Inches(0.35), Inches(5.4), Inches(12.6), Inches(0.8),
    size=13, color=C_LGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TAiCapability enum
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "The TAiCapability Enum",
                "Source: Source/Core/uMakerAi.Core.pas")

add_bullet_box(sl,
    "Input / Understanding  (native via completions)",
    [
        "cap_Image          — model understands incoming images",
        "cap_Audio          — model understands / transcribes audio",
        "cap_Video          — model understands incoming video",
        "cap_Pdf            — model understands PDFs",
        "cap_WebSearch      — native web search",
        "cap_Reasoning      — extended reasoning / CoT / thinking",
        "cap_CodeInterpreter— executes code natively",
        "cap_Memory         — persistent memory",
        "cap_TextEditor     — can edit text files",
        "cap_ComputerUse    — computer control",
        "cap_Shell          — shell command execution",
    ],
    Inches(0.35), Inches(1.2), Inches(6.15), Inches(5.7),
    title_color=C_TEAL, bullet_size=12,
    bullet_char="  • ")

add_bullet_box(sl,
    "Output / Generation  (gap → activates bridge)",
    [
        "cap_GenImage   — produce an image as output",
        "cap_GenAudio   — produce audio output (TTS)",
        "cap_GenVideo   — produce video output",
        "cap_GenReport  — produce PDF / HTML / XLSX report",
        "cap_ExtractCode— extract code blocks from response",
    ],
    Inches(6.85), Inches(1.2), Inches(6.1), Inches(2.85),
    box_color=RGBColor(0x1A, 0x2A, 0x0A),
    title_color=C_GREEN, bullet_size=12,
    bullet_char="  • ")

add_text(sl, "Key rule:",
         Inches(6.85), Inches(4.15), Inches(6.1), Inches(0.35),
         size=13, bold=True, color=C_ORANGE)
add_text(sl,
    "If a Gen* capability is in ModelCaps AND SessionCaps → no gap "
    "→ completions handles it inline.\n\n"
    "If a Gen* capability is ONLY in SessionCaps → gap exists "
    "→ orchestrator redirects to the dedicated endpoint.",
    Inches(6.85), Inches(4.5), Inches(6.1), Inches(2.3),
    size=12, color=C_LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ModelCaps vs SessionCaps
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "ModelCaps vs SessionCaps",
                "The two key properties of TAiChat")

# ModelCaps box
add_rect(sl, Inches(0.35), Inches(1.2), Inches(5.9), Inches(5.7),
         fill=RGBColor(0x0A, 0x1A, 0x35), line=C_BLUE, line_w=Pt(2))
add_text(sl, "ModelCaps", Inches(0.5), Inches(1.3), Inches(5.6), Inches(0.5),
         size=22, bold=True, color=C_TEAL)
add_text(sl, "What the model can do natively",
         Inches(0.5), Inches(1.75), Inches(5.6), Inches(0.35),
         size=13, italic=True, color=C_LGRAY)
add_rect(sl, Inches(0.5), Inches(2.1), Inches(5.6), Inches(0.04), fill=C_BLUE)

entries_m = [
    ("GPT-4.1",        "[cap_Image]"),
    ("dall-e-3",       "[]  (image endpoint only)"),
    ("gemini-2.5-flash","[cap_Image, cap_Audio, cap_Video,\n cap_Pdf, cap_WebSearch,\n cap_Reasoning, cap_CodeInterpreter]"),
    ("deepseek-chat",  "[]"),
    ("deepseek-reasoner","[cap_Reasoning]"),
    ("Claude (all)",   "[cap_Image, cap_Pdf,\n cap_Reasoning, cap_WebSearch]"),
]
ty = Inches(2.2)
for model, caps in entries_m:
    add_text(sl, model, Inches(0.5), ty, Inches(2.1), Inches(0.55),
             size=11, bold=True, color=C_ORANGE)
    add_text(sl, caps, Inches(2.6), ty, Inches(3.1), Inches(0.6),
             size=10, color=C_WHITE)
    ty += Inches(0.58)

# SessionCaps box
add_rect(sl, Inches(6.7), Inches(1.2), Inches(6.25), Inches(5.7),
         fill=RGBColor(0x0A, 0x2A, 0x1A), line=C_GREEN, line_w=Pt(2))
add_text(sl, "SessionCaps", Inches(6.85), Inches(1.3), Inches(6.0), Inches(0.5),
         size=22, bold=True, color=C_GREEN)
add_text(sl, "What this session needs",
         Inches(6.85), Inches(1.75), Inches(6.0), Inches(0.35),
         size=13, italic=True, color=C_LGRAY)
add_rect(sl, Inches(6.85), Inches(2.1), Inches(6.0), Inches(0.04), fill=C_GREEN)

scenarios = [
    ("Use GPT-4o for TTS",
     "ModelCaps = [cap_Image]\nSessionCaps = [cap_Image, cap_GenAudio]\n→ Gap = [cap_GenAudio]\n→ TTS endpoint activated"),
    ("Text-only Ollama + images",
     "ModelCaps = []\nSessionCaps = [cap_Image]\n→ Gap = [cap_Image]\n→ Image-description bridge runs"),
    ("DALL-E 3 image gen",
     "ModelCaps = []\nSessionCaps = [cap_GenImage]\n→ Gap = [cap_GenImage]\n→ Image-generation endpoint"),
    ("Gemini full multimodal",
     "ModelCaps = SessionCaps\n→ Gap = []\n→ Direct completions call"),
]
ty = Inches(2.2)
for title, detail in scenarios:
    add_rect(sl, Inches(6.85), ty, Inches(5.9), Inches(1.18),
             fill=RGBColor(0x0F, 0x32, 0x1A), line=C_GREEN, line_w=Pt(1))
    add_text(sl, title, Inches(7.0), ty + Inches(0.07), Inches(5.6), Inches(0.32),
             size=12, bold=True, color=C_YELLOW)
    add_text(sl, detail, Inches(7.0), ty + Inches(0.37), Inches(5.6), Inches(0.75),
             size=10, color=C_LGRAY)
    ty += Inches(1.28)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Gap Analysis diagram
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Gap Analysis — The Orchestration Engine",
                "TAiChat.RunNew  ·  Source/Core/uMakerAi.Chat.pas:3004")

add_code_box(sl,
    ["Gap := FSessionCaps - FModelCaps;   // set subtraction — runs on every Run() call"],
    Inches(0.35), Inches(1.2), Inches(12.6), Inches(0.55))

# Phase boxes
phases = [
    ("Phase 1", "Input Bridge\n(cmConversation only)",
     ["Gap has cap_Audio  → InternalRunTranscription",
      "Gap has cap_Image  → InternalRunImageDescription",
      "Gap has cap_Pdf    → InternalRunPDFDescription",
      "OnProcessMediaFile event fires first (priority 1)"],
     C_BLUE, RGBColor(0x05, 0x15, 0x40)),
    ("Phase 2", "Grounding\n(always, all modes)",
     ["Gap has cap_WebSearch",
      "→ InternalRunWebSearch",
      "→ results injected into context",
      "before the completions call"],
     C_ORANGE, RGBColor(0x30, 0x18, 0x00)),
    ("Phase 3", "Output Orchestration\n(cmConversation)",
     ["Gap has cap_GenVideo  → InternalRunImageVideoGeneration",
      "Gap has cap_GenImage  → InternalRunImageGeneration",
      "Gap has cap_GenAudio  → InternalRunSpeechGeneration",
      "Gap has cap_GenReport → InternalRunReport",
      "Gap is empty          → InternalRunCompletions"],
     C_GREEN, RGBColor(0x05, 0x25, 0x10)),
]

for i, (ph, title, bullets, col, bg) in enumerate(phases):
    lx = Inches(0.35) + i * Inches(4.35)
    add_rect(sl, lx, Inches(1.9), Inches(4.15), Inches(5.0),
             fill=bg, line=col, line_w=Pt(2))
    add_text(sl, ph, lx + Inches(0.15), Inches(1.98),
             Inches(3.8), Inches(0.38),
             size=17, bold=True, color=col)
    add_text(sl, title, lx + Inches(0.15), Inches(2.34),
             Inches(3.8), Inches(0.55),
             size=12, color=C_LGRAY, italic=True)
    add_rect(sl, lx + Inches(0.15), Inches(2.88), Inches(3.8), Inches(0.03),
             fill=col)
    txb = sl.shapes.add_textbox(lx + Inches(0.15), Inches(2.95),
                                 Inches(3.82), Inches(3.7))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = "▸ " + b
        run.font.size = Pt(12)
        run.font.color.rgb = C_WHITE

# forced-mode note
add_rect(sl, Inches(0.35), Inches(7.0), Inches(12.6), Inches(0.38),
         fill=RGBColor(0x25, 0x18, 0x00), line=C_ORANGE, line_w=Pt(1))
add_text(sl,
    "Forced modes (cmImageGeneration, cmSpeechGeneration, etc.) skip Phases 1 & 3 gap check "
    "and call the target method directly.",
    Inches(0.5), Inches(7.03), Inches(12.3), Inches(0.32),
    size=11, color=C_ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TAiChatFactory configuration
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Configuration via TAiChatFactory",
                "Source/Chat/uMakerAi.Chat.Initializations.pas")

# priority pyramid
add_bullet_box(sl,
    "3-level priority (lowest → highest)",
    [
        "1. Driver defaults  — RegisterDefaultParams in the driver class",
        "2. Global provider defaults  — RegisterUserParam(Driver, Param, Value)",
        "3. Per-model overrides       — RegisterUserParam(Driver, Model, Param, Value)",
    ],
    Inches(0.35), Inches(1.2), Inches(5.9), Inches(1.8),
    title_color=C_ORANGE, bullet_size=12)

add_bullet_box(sl,
    "ThinkingLevel values",
    [
        "tlDefault — let the provider decide",
        "tlLow     — fast, minimal reasoning, lower cost",
        "tlMedium  — balanced (recommended default)",
        "tlHigh    — maximum quality, slower & costlier",
        "Note: ignored if cap_Reasoning is not in ModelCaps",
    ],
    Inches(6.7), Inches(1.2), Inches(6.25), Inches(1.8),
    title_color=C_TEAL, bullet_size=12)

add_text(sl, "Syntax examples:",
         Inches(0.35), Inches(3.1), Inches(12.6), Inches(0.35),
         size=13, bold=True, color=C_LGRAY)

add_code_box(sl, [
    "// Global default for a provider",
    "TAiChatFactory.Instance.RegisterUserParam('OpenAi', 'ModelCaps',   '[cap_Image]');",
    "TAiChatFactory.Instance.RegisterUserParam('OpenAi', 'SessionCaps', '[cap_Image]');",
    "",
    "// Per-model override",
    "TAiChatFactory.Instance.RegisterUserParam('OpenAi', 'o3', 'ModelCaps',    '[cap_Image, cap_Reasoning]');",
    "TAiChatFactory.Instance.RegisterUserParam('OpenAi', 'o3', 'SessionCaps',  '[cap_Image, cap_Reasoning]');",
    "TAiChatFactory.Instance.RegisterUserParam('OpenAi', 'o3', 'ThinkingLevel', 'tlMedium');",
    "",
    "// Custom profile (aa_* prefix convention)",
    "TAiChatFactory.Instance.RegisterCustomModel('OpenAi', 'aa_o3-high', 'o3');",
    "TAiChatFactory.Instance.RegisterUserParam('OpenAi', 'aa_o3-high', 'ThinkingLevel', 'tlHigh');",
], Inches(0.35), Inches(3.5), Inches(12.6), Inches(3.1))

add_text(sl,
    "String format: '[cap_Image, cap_Reasoning]'  — exact enum name, case-sensitive, "
    "space after comma.  Parsed via RTTI (tkSet) in ApplyParamsToChat.",
    Inches(0.35), Inches(6.68), Inches(12.6), Inches(0.45),
    size=11, color=C_LGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Runtime configuration
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Runtime Configuration",
                "Assigning ModelCaps / SessionCaps directly in code")

add_text(sl, "On TAiChatConnection:",
         Inches(0.35), Inches(1.2), Inches(12.6), Inches(0.35),
         size=14, bold=True, color=C_TEAL)
add_code_box(sl, [
    "AiConnection.DriverName := 'OpenAi';",
    "AiConnection.Model      := 'gpt-4.1';",
    "",
    "// Override caps for this specific session",
    "AiConnection.ModelCaps   := [cap_Image, cap_Pdf];",
    "AiConnection.SessionCaps := [cap_Image, cap_Pdf, cap_GenAudio];",
    "// Gap = [cap_GenAudio]  →  next Run() will call the TTS endpoint",
], Inches(0.35), Inches(1.6), Inches(6.1), Inches(2.4))

add_text(sl, "On TAiChat directly:",
         Inches(6.7), Inches(1.2), Inches(6.25), Inches(0.35),
         size=14, bold=True, color=C_TEAL)
add_code_box(sl, [
    "var Chat: TAiOpenChat;",
    "Chat := TAiOpenChat.Create(nil);",
    "Chat.ApiKey := '@OPENAI_API_KEY';",
    "Chat.Model  := 'gpt-image-1';",
    "Chat.ModelCaps   := [];",
    "Chat.SessionCaps := [cap_GenImage];",
    "// Gap=[cap_GenImage] → InternalRunImageGeneration",
], Inches(6.7), Inches(1.6), Inches(6.25), Inches(2.4))

add_text(sl, "Reading the effective caps:",
         Inches(0.35), Inches(4.1), Inches(12.6), Inches(0.35),
         size=14, bold=True, color=C_TEAL)
add_code_box(sl, [
    "var Gap: TAiCapabilities;",
    "Gap := AiConnection.SessionCaps - AiConnection.ModelCaps;",
    "",
    "if cap_GenAudio in Gap then",
    "  ShowMessage('This session will generate audio via TTS');",
    "",
    "if cap_Reasoning in AiConnection.ModelCaps then",
    "  ShowMessage('This model has native extended reasoning');",
], Inches(0.35), Inches(4.5), Inches(12.6), Inches(2.1))

add_rect(sl, Inches(0.35), Inches(6.65), Inches(12.6), Inches(0.65),
         fill=RGBColor(0x25, 0x18, 0x00), line=C_ORANGE, line_w=Pt(1))
add_text(sl,
    "Important: assigning ModelCaps or SessionCaps in code sets FNewSystemConfigured = True, "
    "preventing EnsureNewSystemConfig from overwriting them with an automatic legacy translation.",
    Inches(0.5), Inches(6.72), Inches(12.3), Inches(0.5),
    size=12, color=C_ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Backward compatibility
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Backward Compatibility",
                "Legacy params still work — automatic translation in EnsureNewSystemConfig")

add_bullet_box(sl,
    "How legacy params map to the new system",
    [
        "NativeInputFiles + ChatMediaSupports  →  ModelCaps",
        "NativeOutputFiles + EnabledFeatures   →  SessionCaps",
        "Translation runs once per session, only if FNewSystemConfigured = False",
        "Fully transparent — no code changes required in existing apps",
    ],
    Inches(0.35), Inches(1.2), Inches(12.6), Inches(1.85),
    title_color=C_TEAL, bullet_size=13)

add_text(sl, "Priority chain:",
         Inches(0.35), Inches(3.15), Inches(12.6), Inches(0.35),
         size=14, bold=True, color=C_LGRAY)

priorities = [
    ("Highest", "ModelCaps / SessionCaps assigned explicitly in code",
     "FNewSystemConfigured = True", C_GREEN),
    ("Middle",  "Automatic translation from legacy params",
     "FNewSystemConfigured = False", C_TEAL),
    ("Lowest",  "Driver defaults (RegisterDefaultParams)",
     "class-level defaults", C_ORANGE),
]
for i, (level, desc, detail, col) in enumerate(priorities):
    lx = Inches(0.35) + i * Inches(4.35)
    add_rect(sl, lx, Inches(3.55), Inches(4.15), Inches(1.4),
             fill=C_ACCENT, line=col, line_w=Pt(2))
    add_text(sl, level, lx + Inches(0.15), Inches(3.62),
             Inches(3.8), Inches(0.35), size=14, bold=True, color=col)
    add_text(sl, desc, lx + Inches(0.15), Inches(3.97),
             Inches(3.8), Inches(0.5), size=11, color=C_WHITE)
    add_text(sl, detail, lx + Inches(0.15), Inches(4.47),
             Inches(3.8), Inches(0.35), size=10, color=C_LGRAY, italic=True)

add_text(sl, "Migration — before vs after:",
         Inches(0.35), Inches(5.1), Inches(12.6), Inches(0.35),
         size=14, bold=True, color=C_LGRAY)
add_code_box(sl, [
    "// Before (legacy)",
    "TAiChatFactory.Instance.RegisterUserParam('MyDriver', 'ChatMediaSupports', 'Tcm_Image,Tcm_Pdf');",
    "TAiChatFactory.Instance.RegisterUserParam('MyDriver', 'EnabledFeatures',   'Tcm_Image,Tcm_Pdf');",
    "TAiChatFactory.Instance.RegisterUserParam('MyDriver', 'NativeInputFiles',  'Tfc_Image,Tfc_Pdf');",
    "",
    "// After (v3.3 — same result, much simpler)",
    "TAiChatFactory.Instance.RegisterUserParam('MyDriver', 'ModelCaps',   '[cap_Image, cap_Pdf]');",
    "TAiChatFactory.Instance.RegisterUserParam('MyDriver', 'SessionCaps', '[cap_Image, cap_Pdf]');",
], Inches(0.35), Inches(5.5), Inches(12.6), Inches(1.78))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Configuration Patterns (part 1: 1-5)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Common Configuration Patterns  (1 – 5)",
                "RegisterUserParam shorthand used for brevity")

patterns = [
    ("1  Plain text model",
     ["ModelCaps   = '[]'",
      "SessionCaps = '[]'",
      "Tool_Active = 'False'"],
     "Gap = []  →  direct InternalRunCompletions"),
    ("2  Native vision",
     ["ModelCaps   = '[cap_Image]'",
      "SessionCaps = '[cap_Image]'",
      "Tool_Active = 'True'"],
     "Gap = []  →  images go straight to completions API"),
    ("3  Full multimodal  (Gemini 2.5 Flash)",
     ["ModelCaps = SessionCaps =",
      "[cap_Image, cap_Audio, cap_Video,",
      " cap_Pdf, cap_WebSearch,",
      " cap_Reasoning, cap_CodeInterpreter]"],
     "Gap = []  →  everything native in completions"),
    ("4  Reasoning model  (o3, deepseek-reasoner)",
     ["ModelCaps   = '[cap_Image, cap_Reasoning]'",
      "SessionCaps = '[cap_Image, cap_Reasoning]'",
      "ThinkingLevel = 'tlMedium'"],
     "Gap = []  →  driver activates extended thinking"),
    ("5  TTS via dedicated endpoint",
     ["ModelCaps   = '[]'",
      "SessionCaps = '[cap_GenAudio]'",
      "Tool_Active = 'False'",
      "Voice       = 'alloy'"],
     "Gap = [cap_GenAudio]  →  InternalRunSpeechGeneration"),
]

for i, (title, code, result) in enumerate(patterns):
    col = i % 2
    row = i // 2
    if i == 4:   # 5th card spans bottom centered
        lx = Inches(0.35) + 1 * Inches(6.45)
        ty = Inches(1.2) + 2 * Inches(2.35)
    else:
        lx = Inches(0.35) + col * Inches(6.45)
        ty = Inches(1.2) + row * Inches(2.35)
    add_rect(sl, lx, ty, Inches(6.15), Inches(2.2),
             fill=C_ACCENT, line=C_BLUE, line_w=Pt(1))
    add_text(sl, title, lx + Inches(0.12), ty + Inches(0.08),
             Inches(5.9), Inches(0.38), size=13, bold=True, color=C_TEAL)
    add_code_box(sl, code,
                 lx + Inches(0.12), ty + Inches(0.5),
                 Inches(5.9), Inches(1.0), size=10)
    add_text(sl, result,
             lx + Inches(0.12), ty + Inches(1.54),
             Inches(5.9), Inches(0.5),
             size=11, color=C_GREEN, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Configuration Patterns (part 2: 6-10)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Common Configuration Patterns  (6 – 10)",
                "RegisterUserParam shorthand used for brevity")

patterns2 = [
    ("6  Image gen — dedicated endpoint  (DALL-E)",
     ["ModelCaps   = '[]'",
      "SessionCaps = '[cap_GenImage]'",
      "Tool_Active = 'False'"],
     "Gap=[cap_GenImage]  →  InternalRunImageGeneration"),
    ("7  Image gen — native in completions  (Gemini)",
     ["ModelCaps   = '[cap_Image, cap_GenImage]'",
      "SessionCaps = '[cap_Image, cap_GenImage]'",
      "Tool_Active = 'False'"],
     "Gap = []  →  completions returns image inline"),
    ("8  STT — transcription endpoint  (Whisper)",
     ["ModelCaps   = '[cap_Audio]'",
      "SessionCaps = '[cap_Audio]'",
      "Tool_Active = 'False'",
      "ChatMode    = cmTranscription"],
     "Gap = []  →  InternalRunTranscription"),
    ("9  Vision bridge for text-only model  (Ollama)",
     ["ModelCaps   = '[]'   // no native vision",
      "SessionCaps = '[cap_Image]'"],
     "Gap=[cap_Image]  →  Phase 1 auto-describes images in text"),
    ("10  Video generation  (Veo / Grok)",
     ["ModelCaps   = '[cap_Image]'",
      "SessionCaps = '[cap_Image, cap_GenVideo]'",
      "Tool_Active = 'False'"],
     "Gap=[cap_GenVideo]  →  InternalRunImageVideoGeneration"),
]

for i, (title, code, result) in enumerate(patterns2):
    col = i % 2
    row = i // 2
    if i == 4:
        lx = Inches(0.35) + 1 * Inches(6.45)
        ty = Inches(1.2) + 2 * Inches(2.35)
    else:
        lx = Inches(0.35) + col * Inches(6.45)
        ty = Inches(1.2) + row * Inches(2.35)
    add_rect(sl, lx, ty, Inches(6.15), Inches(2.2),
             fill=C_ACCENT, line=C_BLUE, line_w=Pt(1))
    add_text(sl, title, lx + Inches(0.12), ty + Inches(0.08),
             Inches(5.9), Inches(0.38), size=13, bold=True, color=C_TEAL)
    add_code_box(sl, code,
                 lx + Inches(0.12), ty + Inches(0.5),
                 Inches(5.9), Inches(1.0), size=10)
    add_text(sl, result,
             lx + Inches(0.12), ty + Inches(1.54),
             Inches(5.9), Inches(0.5),
             size=11, color=C_GREEN, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Provider Reference
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Provider Capability Reference",
                "Source: uMakerAi.Chat.Initializations.pas  ·  Feb 2026")

rows = [
    # provider, model/global, ModelCaps abbrev, Gap note, Tool_Active
    ("OpenAI",   "gpt-4.1 / 4.1-mini",        "[cap_Image]",                 "—",                    "✔"),
    ("OpenAI",   "o3 / o4-mini",               "[cap_Image, cap_Reasoning]",  "— (tlMedium)",         "✔"),
    ("OpenAI",   "dall-e-3 / gpt-image-1",     "[]",                          "→ Image endpoint",     "✘"),
    ("OpenAI",   "gpt-4o-mini-tts",            "[]",                          "→ TTS endpoint",       "✘"),
    ("OpenAI",   "gpt-4o-audio-preview",       "[cap_Audio, cap_GenAudio]",   "—",                    "✘"),
    ("Gemini",   "gemini-2.5-flash",           "[Image,Audio,Video,Pdf,\nWebSearch,Reasoning,Code]", "—", "✔"),
    ("Gemini",   "gemini-2.5-flash-image",     "[cap_Image, cap_GenImage]",   "—",                    "✘"),
    ("Gemini",   "aa_veo-*",                   "[cap_Image]",                 "→ Video endpoint",     "✘"),
    ("Claude",   "All current models",         "[Image,Pdf,Reasoning,\nWebSearch]", "—",             "✘"),
    ("Groq",     "compound-beta",              "[WebSearch, CodeInterpreter]","—",                    "✘"),
    ("Groq",     "whisper-large-v3",           "[cap_Audio]",                 "—",                    "✘"),
    ("Groq",     "Orpheus TTS",                "[]",                          "→ TTS endpoint",       "✘"),
    ("DeepSeek", "deepseek-reasoner",          "[cap_Reasoning]",             "— (tlMedium)",         "✔"),
    ("Mistral",  "magistral-medium/small",     "[cap_Reasoning]",             "— (tlMedium)",         "✔"),
    ("Ollama",   "qwen3:latest",               "[cap_Reasoning]",             "— (tlMedium)",         "✔"),
    ("Ollama",   "gemma3 / qwen2.5vl",         "[cap_Image]",                 "—",                    "✔"),
]

headers = ["Provider", "Model / Global", "ModelCaps", "Gap / Note", "Tools"]
col_widths = [Inches(1.35), Inches(2.35), Inches(4.35), Inches(2.45), Inches(0.9)]
col_starts = [Inches(0.25)]
for cw in col_widths[:-1]:
    col_starts.append(col_starts[-1] + cw)

row_h = Inches(0.355)
header_y = Inches(1.2)

# header row
add_rect(sl, Inches(0.25), header_y, sum(col_widths), row_h, fill=C_BLUE)
for j, (hdr, cx, cw) in enumerate(zip(headers, col_starts, col_widths)):
    add_text(sl, hdr, cx + Inches(0.05), header_y + Inches(0.05),
             cw - Inches(0.1), row_h - Inches(0.08),
             size=11, bold=True, color=C_WHITE)

for i, (prov, model, mcaps, gap, tools) in enumerate(rows):
    ry = header_y + (i + 1) * row_h
    bg = C_ACCENT if i % 2 == 0 else RGBColor(0x12, 0x12, 0x28)
    add_rect(sl, Inches(0.25), ry, sum(col_widths), row_h, fill=bg)
    tool_col = C_GREEN if tools == "✔" else RGBColor(0xAA, 0x44, 0x44)
    for j, (val, cx, cw) in enumerate(zip(
            [prov, model, mcaps, gap, tools], col_starts, col_widths)):
        tc = C_TEAL if j == 0 else (tool_col if j == 4 else C_WHITE)
        add_text(sl, val, cx + Inches(0.05), ry + Inches(0.04),
                 cw - Inches(0.1), row_h - Inches(0.06),
                 size=9, color=tc, bold=(j == 0))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Adding a new provider
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Adding a New Provider",
                "5 design questions + configuration template")

add_bullet_box(sl,
    "Design checklist",
    [
        "1. What input types does the completions endpoint accept?  →  set ModelCaps (cap_Image, cap_Audio…)",
        "2. Can completions return images/audio/video inline?  →  add cap_Gen* to BOTH ModelCaps & SessionCaps",
        "3. Does the model use a separate TTS/image/video endpoint?  →  ModelCaps=[], SessionCaps=[cap_Gen*]",
        "4. Does the model support extended reasoning?  →  cap_Reasoning + ThinkingLevel",
        "5. Does the model support tool calling?  →  Tool_Active = True",
    ],
    Inches(0.35), Inches(1.2), Inches(12.6), Inches(2.1),
    title_color=C_ORANGE, bullet_size=12)

add_text(sl, "Configuration template:",
         Inches(0.35), Inches(3.4), Inches(12.6), Inches(0.35),
         size=14, bold=True, color=C_LGRAY)

add_code_box(sl, [
    "// 1. Global provider defaults",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'Max_Tokens',  '16000');",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'Tool_Active', 'True');",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'ModelCaps',   '[cap_Image]');",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'SessionCaps', '[cap_Image]');",
    "",
    "// 2. Text-only model override",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'my-text-model', 'ModelCaps',   '[]');",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'my-text-model', 'SessionCaps', '[]');",
    "",
    "// 3. Reasoning model",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'my-think-model', 'ModelCaps',    '[cap_Image, cap_Reasoning]');",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'my-think-model', 'SessionCaps',  '[cap_Image, cap_Reasoning]');",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'my-think-model', 'ThinkingLevel', 'tlMedium');",
    "",
    "// 4. TTS model",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'my-tts-model', 'ModelCaps',   '[]');",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'my-tts-model', 'SessionCaps', '[cap_GenAudio]');",
    "TAiChatFactory.Instance.RegisterUserParam('MyProvider', 'my-tts-model', 'Tool_Active', 'False');",
], Inches(0.35), Inches(3.78), Inches(12.6), Inches(3.3))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Summary & Source reference
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_DARK)
slide_title_bar(sl, "Summary & Source File Reference", "")

add_bullet_box(sl,
    "Key takeaways",
    [
        "ModelCaps = what the model does natively via completions",
        "SessionCaps = what the session needs (may exceed native caps)",
        "Gap = SessionCaps − ModelCaps  →  auto-activates the right bridge",
        "Configure once in Initializations.pas; override at runtime if needed",
        "Fully backward-compatible — existing legacy configs keep working",
        "10 patterns cover every real-world scenario",
    ],
    Inches(0.35), Inches(1.2), Inches(7.0), Inches(3.8),
    title_color=C_TEAL, bullet_size=13)

# source file table
src = [
    ("TAiCapability enum definition",              "Source/Core/uMakerAi.Core.pas : 65"),
    ("ModelCaps / SessionCaps properties",         "Source/Core/uMakerAi.Chat.pas : 430"),
    ("Setter + legacy sync implementation",        "Source/Core/uMakerAi.Chat.pas : 2886"),
    ("Gap analysis + RunNew orchestration",        "Source/Core/uMakerAi.Chat.pas : 3004"),
    ("RTTI param application (ApplyParamsToChat)", "Source/Chat/uMakerAi.Chat.AiConnection.pas : 509"),
    ("Per-provider model configuration",           "Source/Chat/uMakerAi.Chat.Initializations.pas"),
]
add_rect(sl, Inches(7.7), Inches(1.2), Inches(5.25), Inches(3.8),
         fill=C_ACCENT, line=C_BLUE, line_w=Pt(1))
add_text(sl, "Source file reference",
         Inches(7.85), Inches(1.28), Inches(5.0), Inches(0.38),
         size=14, bold=True, color=C_ORANGE)
for i, (purpose, loc) in enumerate(src):
    ry = Inches(1.7) + i * Inches(0.52)
    add_text(sl, purpose, Inches(7.85), ry, Inches(5.0), Inches(0.28),
             size=10, color=C_WHITE)
    add_text(sl, loc, Inches(7.85), ry + Inches(0.27), Inches(5.0), Inches(0.22),
             size=9, color=C_TEAL, italic=True)

# closing strip
add_rect(sl, 0, SLIDE_H - Inches(1.25), SLIDE_W, Inches(1.25), fill=C_BLUE)
add_text(sl, "MakerAI v3.3  —  Capability System",
         Inches(0.4), SLIDE_H - Inches(1.15), Inches(8), Inches(0.45),
         size=20, bold=True, color=C_WHITE)
add_text(sl, "makerai.cimamaker.com",
         Inches(0.4), SLIDE_H - Inches(0.72), Inches(8), Inches(0.38),
         size=14, color=C_TEAL)
add_text(sl, "March 2026",
         Inches(10.5), SLIDE_H - Inches(0.8), Inches(2.5), Inches(0.38),
         size=13, color=C_LGRAY, align=PP_ALIGN.RIGHT)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = r"E:\Delphi\Delphi13\Compo\FMXCompo\AiMaker\Docs\Version 3\uMakerAi-CapabilitySystem.EN.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
